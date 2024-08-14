"""
Document Processing Module

This module provides functionality to load various document types, process them,
and create a vector store for embedding-based search. The module supports PDF,
DOCX, TXT, CSV, PPTX, and Excel files, with easy extensibility to support other formats.
"""

import tempfile
import os
import logging
from typing import Dict, Any
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.document_loaders import PyPDFLoader, TextLoader, CSVLoader
import pandas as pd
from docx import Document
from pptx import Presentation
import asyncio
from sqlalchemy import create_engine, Column, Integer, String, LargeBinary
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker
import hashlib

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

SUPPORTED_FILE_TYPES = {
    "pdf": 10 * 1024 * 1024,  
    "docx": 5 * 1024 * 1024,  
    "txt": 1 * 1024 * 1024,   
    "csv": 5 * 1024 * 1024,   
    "xlsx": 10 * 1024 * 1024, 
    "xls": 10 * 1024 * 1024,  
    "pptx": 20 * 1024 * 1024  
}

CONFIG = {
    "chunk_size": 1000,
    "chunk_overlap": 200,
    "embedding_model": "sentence-transformers/all-MiniLM-L6-v2",
    "db_url": "sqlite:///documents.db"
}

Base = declarative_base()

class Document(Base):
    __tablename__ = 'documents'

    id = Column(Integer, primary_key=True)
    filename = Column(String)
    file_type = Column(String)
    content_hash = Column(String)
    vector_store = Column(LargeBinary)

engine = create_engine(CONFIG["db_url"])
Base.metadata.create_all(engine)
Session = sessionmaker(bind=engine)

def validate_file(file, file_type: str) -> None:
    if file_type not in SUPPORTED_FILE_TYPES:
        raise ValueError(f"Unsupported file type: {file_type}")
    
    file_size = file.getbuffer().nbytes
    if file_size > SUPPORTED_FILE_TYPES[file_type]:
        raise ValueError(f"File too large. Maximum size for {file_type} is {SUPPORTED_FILE_TYPES[file_type] / (1024 * 1024)} MB")

def load_docx(file_path: str) -> str:
    try:
        doc = Document(file_path)
        text = [para.text for para in doc.paragraphs]
        if not text:
            raise ValueError("DOCX file is empty or cannot be read.")
        return '\n'.join(text)
    except Exception as e:
        logging.error(f"Error loading DOCX file: {e}")
        raise ValueError(f"Error loading DOCX file: {e}")

def load_excel(file_path: str) -> str:
    try:
        df = pd.read_excel(file_path, engine='openpyxl')
        if df.empty:
            raise ValueError("Excel file is empty or cannot be read.")
        return df.to_csv(index=False)
    except Exception as e:
        logging.error(f"Error loading Excel file: {e}")
        raise ValueError(f"Error loading Excel file: {e}")

def load_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        if not text_runs:
            raise ValueError("PowerPoint file is empty or cannot be read.")
        return '\n'.join(text_runs)
    except Exception as e:
        logging.error(f"Error loading PowerPoint file: {e}")
        raise ValueError(f"Error loading PowerPoint file: {e}")

def load_document(file_path: str, file_type: str) -> list:
    try:
        if file_type == "pdf":
            loader = PyPDFLoader(file_path)
            documents = loader.load()
        elif file_type == "docx":
            text = load_docx(file_path)
            documents = [{'text': text, 'metadata': {'source': file_path}}]
        elif file_type == "txt":
            loader = TextLoader(file_path)
            documents = loader.load()
        elif file_type == "csv":
            loader = CSVLoader(file_path)
            documents = loader.load()
        elif file_type in ["xlsx", "xls"]:
            text = load_excel(file_path)
            documents = [{'text': text, 'metadata': {'source': file_path}}]
        elif file_type == "pptx":
            text = load_pptx(file_path)
            documents = [{'text': text, 'metadata': {'source': file_path}}]
        else:
            raise ValueError("Unsupported file type")

        for doc in documents:
            if 'metadata' not in doc:
                doc['metadata'] = {'source': file_path}
            else:
                doc['metadata']['source'] = file_path

        return documents
    except Exception as e:
        logging.error(f"Error loading document: {e}")
        raise ValueError(f"Error loading document: {e}")

async def process_document(file_path: str, file_type: str) -> FAISS:
    documents = await asyncio.to_thread(load_document, file_path, file_type)
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=CONFIG["chunk_size"],
        chunk_overlap=CONFIG["chunk_overlap"]
    )
    
    texts_with_metadata = []
    for doc in documents:
        text_chunks = text_splitter.split_text(doc['text'])
        for chunk in text_chunks:
            texts_with_metadata.append({
                'text': chunk,
                'metadata': doc['metadata']
            })

    embeddings = HuggingFaceEmbeddings(model_name=CONFIG["embedding_model"])
    vectorstore = await asyncio.to_thread(FAISS.from_documents, texts_with_metadata, embeddings)
    
    return vectorstore

async def load_document_and_create_vectorstore(uploaded_file: Any, file_type: str) -> FAISS:
    validate_file(uploaded_file, file_type)

    temp_file_path = None

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}", mode='wb') as temp_file:
            temp_file.write(uploaded_file.read())
            temp_file_path = temp_file.name
            logging.info(f"Temporary file created at: {temp_file_path}")

        if not os.path.exists(temp_file_path):
            raise FileNotFoundError(f"Temporary file not found: {temp_file_path}")

        file_hash = hashlib.sha256(open(temp_file_path, 'rb').read()).hexdigest()

        session = Session()
        existing_doc = session.query(Document).filter_by(content_hash=file_hash).first()
        if existing_doc:
            logging.info(f"Document with hash {file_hash} already exists in database. Retrieving existing vector store.")
            return FAISS.deserialize_from_bytes(existing_doc.vector_store)

        vectorstore = await process_document(temp_file_path, file_type)

        new_doc = Document(
            filename=uploaded_file.name,
            file_type=file_type,
            content_hash=file_hash,
            vector_store=vectorstore.serialize_to_bytes()
        )
        session.add(new_doc)
        session.commit()

        return vectorstore

    except Exception as e:
        logging.error(f"Error processing document and creating vector store: {e}")
        raise RuntimeError(f"Error processing document and creating vector store: {e}")

    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
            logging.info(f"Temporary file deleted: {temp_file_path}")
        if 'session' in locals():
            session.close()

async def process_uploaded_file(uploaded_file):
    file_type = uploaded_file.name.split('.')[-1].lower()
    try:
        vectorstore = await load_document_and_create_vectorstore(uploaded_file, file_type)
        # Use the vectorstore for further processing or querying
    except Exception as e:
        logging.error(f"Error processing uploaded file: {e}")
